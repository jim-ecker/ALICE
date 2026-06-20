"""PDF metadata extraction and LLM-based inference for local folder ingestion."""
from __future__ import annotations

import json
import re
from pathlib import Path

# PDF /Title values that are tool artifacts, not real titles.
_BAD_TITLE_PATTERNS = re.compile(
    r"microsoft\s+word|openoffice|libreoffice|latex|untitled|\.docx?|\.tex\b",
    re.IGNORECASE,
)

# Lines that signal we've moved past the title into author/affiliation blocks.
_AUTHOR_BLOCK_PATTERNS = re.compile(
    r"@|\buniversity\b|\binstitute\b|\blaborator\b|\bdepartment\b|"
    r"\bcenter\b|\bnasa\b|\babstract\b|\bintroduction\b|\breceived\b",
    re.IGNORECASE,
)


def _title_from_first_page(text: str) -> str:
    """Heuristically extract a document title from first-page text.

    Research paper titles appear as the first substantial text block before
    author names and affiliations. This function:
      - Skips blank lines and very short lines (page numbers, headers)
      - Collects lines that look like a title (≥15 chars, no author signals)
      - Stops at the first line that looks like an author/affiliation block
      - Returns at most 3 joined lines (titles rarely span more)
    """
    title_lines: list[str] = []
    for raw in text.splitlines():
        line = raw.strip()
        if not line:
            # A blank line after we've started collecting means the title block ended
            if title_lines:
                break
            continue
        if len(line) < 8:
            # Skip page numbers, section labels, lone characters
            if title_lines:
                break
            continue
        if _AUTHOR_BLOCK_PATTERNS.search(line):
            break
        title_lines.append(line)
        if len(title_lines) == 3:
            break

    return " ".join(title_lines).strip()


def extract_pdf_metadata(pdf_path: Path) -> dict:
    """Extract metadata from a PDF's /Info header dictionary and first-page text.

    Returns a dict with keys: title, authors, year, abstract.
    Values are empty strings where not found.
    """
    import pypdf

    meta: dict[str, str] = {"title": "", "authors": "", "year": "", "abstract": ""}
    first_page_text = ""

    try:
        reader = pypdf.PdfReader(str(pdf_path))
        info = reader.metadata or {}

        raw_title = info.get("/Title", "") or ""
        raw_author = info.get("/Author", "") or ""
        raw_subject = info.get("/Subject", "") or ""
        raw_date = info.get("/CreationDate", "") or ""

        # Use the PDF header title only if it looks genuine
        if raw_title.strip() and not _BAD_TITLE_PATTERNS.search(raw_title):
            meta["title"] = raw_title.strip()
        meta["authors"] = raw_author.strip()
        meta["abstract"] = raw_subject.strip()

        # Extract year from CreationDate string (e.g. "D:20210315..." or "2021-03-15")
        if raw_date:
            year_match = re.search(r"(\d{4})", raw_date)
            if year_match:
                meta["year"] = year_match.group(1)

        # Extract first-page text for title heuristic and LLM inference
        if reader.pages:
            first_page_text = reader.pages[0].extract_text() or ""

        # Fall back to first-page heuristic if header title was absent or bad
        if not meta["title"] and first_page_text:
            meta["title"] = _title_from_first_page(first_page_text)

    except Exception:
        pass

    meta["_first_page_text"] = first_page_text[:3000]
    return meta


def extract_pptx_metadata(path: Path) -> dict:
    """Extract metadata from a PowerPoint file.

    Priority for title:
      1. Title placeholder on slide 1 (most reliable for decks)
      2. Document core properties /Title
      3. First non-empty text on slide 1 (body text fallback)
    """
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN

    meta: dict[str, str] = {"title": "", "authors": "", "year": "", "abstract": ""}
    first_slide_text = ""

    try:
        prs = Presentation(str(path))
        props = prs.core_properties

        # Document properties
        raw_author = (props.author or "").strip()
        meta["authors"] = raw_author
        if props.created:
            meta["year"] = str(props.created.year)
        props_title = (props.title or "").strip()
        if props_title and not _BAD_TITLE_PATTERNS.search(props_title):
            meta["title"] = props_title

        # Slide 1: prefer the title placeholder, fall back to first text block
        if prs.slides:
            slide = prs.slides[0]
            slide_texts: list[str] = []
            title_shape_text = ""
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                # Title placeholder (placeholder type 0 or 1 = TITLE / CENTER_TITLE)
                if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                    ph_idx = shape.placeholder_format.idx
                    if ph_idx in (0, 1) and not title_shape_text:
                        title_shape_text = text
                slide_texts.append(text)

            if title_shape_text:
                meta["title"] = title_shape_text
            elif not meta["title"] and slide_texts:
                # No title placeholder found — use the first substantial text block
                for t in slide_texts:
                    if len(t) >= 8:
                        meta["title"] = t
                        break

            first_slide_text = "\n".join(slide_texts)

    except Exception:
        pass

    meta["_first_page_text"] = first_slide_text[:3000]
    return meta


def extract_excel_metadata(path: Path) -> dict:
    """Extract metadata from an Excel workbook's built-in properties."""
    import openpyxl
    meta: dict[str, str] = {"title": "", "authors": "", "year": "", "abstract": ""}
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        props = wb.properties
        meta["title"] = props.title or ""
        meta["authors"] = props.creator or ""
        if props.created:
            meta["year"] = str(props.created.year)
        wb.close()
    except Exception:
        pass
    if not meta["title"]:
        meta["title"] = path.stem
    return meta


def infer_missing_metadata(partial_meta: dict, llm) -> dict:
    """Use the LLM to fill in missing metadata fields from first-page text.

    Only called when title or authors is absent. Returns updated dict.
    Fields inferred by LLM are marked with a '_inferred' set in the result.
    """
    first_page_text = partial_meta.get("_first_page_text", "")
    if not first_page_text:
        return partial_meta

    missing = [k for k in ("title", "authors", "year", "abstract") if not partial_meta.get(k)]
    if not missing:
        return partial_meta

    prompt = (
        f"You are a research librarian. Based only on the text below (the first page of a PDF), "
        f"extract the following fields: {', '.join(missing)}.\n\n"
        f"Return ONLY a JSON object with those keys. Use empty string if you cannot determine a value.\n\n"
        f"TEXT:\n{first_page_text}"
    )
    messages = [
        {"role": "system", "content": "You extract bibliographic metadata from academic document text. Return only valid JSON."},
        {"role": "user", "content": prompt},
    ]

    inferred_fields: set[str] = set()
    try:
        response = llm.chat(messages, 400)
        # Extract JSON from response (may be wrapped in markdown code block)
        json_match = re.search(r"\{[^{}]+\}", response, re.DOTALL)
        if json_match:
            data = json.loads(json_match.group())
            for key in missing:
                value = (data.get(key) or "").strip()
                if value:
                    partial_meta[key] = value
                    inferred_fields.add(key)
    except Exception:
        pass

    partial_meta["_inferred_fields"] = inferred_fields
    return partial_meta
