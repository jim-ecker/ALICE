from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = "/Users/jecker/Documents/Presentations/Paper Presentation 22 Oct.pptx"
print("Listing shape types for {}...".format(pptx_path))

prs = Presentation(pptx_path)

for slide_num, slide in enumerate(prs.slides, start=1):
    for shape in slide.shapes:
        print("Slide {}, Shape Type: {}".format(slide_num, shape.shape_type))
