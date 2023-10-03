import os
from typing import List
import click
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_text_from_pptx(pptx_path: str) -> List[str]:
    """
    Extract text from a PowerPoint slide deck.

    :param pptx_path: str, path to the PowerPoint file
    :return: list, a list of extracted text strings
    """
    prs = Presentation(pptx_path)
    text_list = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_list.append(f"Slide {slide_num}: {shape.text}")

    return text_list


def extract_images_and_charts_from_pptx(pptx_path: str, output_folder: str):
    """
    Extract images and charts from a PowerPoint slide deck.

    :param pptx_path: str, path to the PowerPoint file
    :param output_folder: str, path to the folder where the extracted images and charts will be saved
    """
    prs = Presentation(pptx_path)
    image_count = 0
    chart_count = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count = extract_images_from_shape(
                    shape, slide_num, output_folder, image_count
                )
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:  # Update this line
                chart_count = extract_charts_from_shape(
                    shape, slide_num, output_folder, chart_count
                )

    click.echo(
        f"Images and charts extracted from {pptx_path} and saved to {output_folder}"
    )


def extract_images_from_shape(shape, slide_num, output_folder, image_count):
    """
    Recursively extract images from a shape and its child shapes.

    :param shape: a PowerPoint shape object
    :param slide_num: int, the slide number that the shape belongs to
    :param output_folder: str, path to the folder where the extracted images will be saved
    :param image_count: int, the current image count
    :return: int, the updated image count
    """
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        parent_folder = os.path.join(
            output_folder, "images"
        )  # Add this line to create the parent "images" folder
        os.makedirs(
            parent_folder, exist_ok=True
        )  # Create the parent "images" folder if it doesn't exist

        slide_folder = os.path.join(
            parent_folder, f"slide_{slide_num}"
        )  # Update this line to use the parent_folder
        os.makedirs(slide_folder, exist_ok=True)

        image = shape.image
        image_filename = f"image_{image_count}.{image.ext}"
        image_output_path = os.path.join(slide_folder, image_filename)

        with open(image_output_path, "wb") as output_file:
            output_file.write(image.blob)

        image_count += 1
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            image_count = extract_images_from_shape(
                child_shape, slide_num, output_folder, image_count
            )

    return image_count


def extract_charts_from_shape(shape, slide_num, output_folder, chart_count):
    """
    Recursively extract charts from a shape and its child shapes.

    :param shape: a PowerPoint shape object
    :param slide_num: int, the slide number that the shape belongs to
    :param output_folder: str, path to the folder where the extracted charts will be saved
    :param chart_count: int, the current chart count
    :return: int, the updated chart count
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
        slide_folder = os.path.join(output_folder, f"slide_{slide_num}")
        os.makedirs(slide_folder, exist_ok=True)

        chart = shape.chart
        chart_filename = f"chart_{chart_count}.{chart.chart_type.ext}"
        chart_output_path = os.path.join(slide_folder, chart_filename)

        chart.save(chart_output_path)

        chart_count += 1
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            chart_count = extract_charts_from_shape(
                child_shape, slide_num, output_folder, chart_count
            )
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        for row in shape.rows:
            for cell in row.cells:
                for element in cell:
                    chart_count = extract_charts_from_shape(
                        element, slide_num, output_folder, chart_count
                    )
    elif shape.shape_type == MSO_SHAPE_TYPE.GRAPHIC_FRAME:
        chart = shape.chart
        if chart is not None:
            slide_folder = os.path.join(output_folder, f"slide_{slide_num}")
            os.makedirs(slide_folder, exist_ok=True)

            chart_filename = f"chart_{chart_count}.{chart.chart_type.ext}"
            chart_output_path = os.path.join(slide_folder, chart_filename)

            chart.save(chart_output_path)

            chart_count += 1

    return chart_count


# ... (existing functions: extract_text_from_pptx, extract_images_and_charts_from_pptx, extract_images_from_shape, extract_charts_from_shape)


@click.command()
@click.option(
    "--pptx",
    type=click.Path(exists=True, file_okay=True, dir_okay=False, readable=True),
    required=True,
    help="Path to the PowerPoint file.",
)
@click.option(
    "--output",
    type=click.Path(file_okay=False, dir_okay=True, writable=True),
    default="./output",
    help="Path to the output folder where the extracted images and charts will be saved.",
)
@click.option(
    "--text",
    is_flag=True,
    default=False,
    help="Enable to extract text from the PowerPoint slides.",
)
def cli(pptx: str, output: str, text: bool):
    """
    Extract text, images, and charts from a PowerPoint presentation.

    :param pptx: str, path to the PowerPoint file
    :param output: str, path to the folder where the extracted images and charts will be saved
    :param text: bool, set to True to extract text from the PowerPoint slides
    """
    os.makedirs(
        output, exist_ok=True
    )  # Add this line to create the output folder if it doesn't exist

    if text:
        text_list = extract_text_from_pptx(pptx)
        with open(os.path.join(output, "text_output.txt"), "w") as text_file:
            text_file.write("\n".join(text_list))
        click.echo(
            f"Text extracted from {pptx} and saved to {os.path.join(output, 'text_output.txt')}"
        )

    extract_images_and_charts_from_pptx(pptx, output)


if __name__ == "__main__":
    cli()
